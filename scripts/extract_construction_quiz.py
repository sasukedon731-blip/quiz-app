from __future__ import annotations

import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


Q_RE = re.compile(r"Q\s*(\d+)")
CHOICE_RE = re.compile(r"^[A-DＡ-Ｄ][\s　]+(.+)$")

LETTER_TO_INDEX = {
    "A": 0,
    "B": 1,
    "C": 2,
    "D": 3,
    "Ａ": 0,
    "Ｂ": 1,
    "Ｃ": 2,
    "Ｄ": 3,
}


def normalize_text(s: str) -> str:
    s = s.replace("\u3000", " ")
    s = s.replace("\xa0", " ")
    s = re.sub(r"[ \t]+", " ", s)
    return s.strip()


def clean_question_text(text: str) -> str:
    text = text.replace("🔍 写真問題", "")
    text = text.replace("🖼 イラスト問題", "")
    text = text.replace("📖 解説文問題", "")
    text = text.replace("📝 道具名問題", "")
    text = text.replace("写真問題", "")
    text = text.replace("イラスト問題", "")
    text = text.replace("解説文問題", "")
    text = text.replace("道具名問題", "")

    # 共通見出しを除去
    text = re.sub(r".*クイズ\s*基本編", "", text)

    lines = [normalize_text(line) for line in text.splitlines()]
    lines = [line for line in lines if line]
    return "\n".join(lines).strip()


def clean_explanation_text(text: str) -> str:
    text = text.replace("📋 解説", "")

    # 「解説」単独見出しは消すが、本文内の語は壊さない
    lines = [normalize_text(line) for line in text.splitlines()]
    lines = [line for line in lines if line and line != "解説"]

    text = "\n".join(lines)
    text = re.sub(r".*クイズ\s*基本編", "", text)

    lines = [normalize_text(line) for line in text.splitlines()]
    lines = [line for line in lines if line]
    return "\n".join(lines).strip()


def clean_hint_text(text: str | None) -> str | None:
    if not text:
        return None

    text = normalize_text(text)
    text = text.replace("ヒント :", "ヒント:")
    text = text.replace("ヒント：", "ヒント:")
    return text.strip()


def get_slide_lines(slide) -> list[str]:
    lines: list[str] = []

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue

        tf = shape.text_frame
        for para in tf.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            text = normalize_text(text)
            if text:
                lines.append(text)

    return lines


def find_q_number(lines: list[str]) -> int | None:
    for line in lines:
        m = Q_RE.search(line)
        if m:
            return int(m.group(1))
    return None


def is_answer_list_slide(lines: list[str]) -> bool:
    return any("答え合わせ" in line for line in lines)


def is_explanation_slide(lines: list[str]) -> bool:
    return any("正解・解説" in line for line in lines)


def detect_kind(lines: list[str]) -> str:
    joined = "\n".join(lines)

    if "イラスト問題" in joined or "写真問題" in joined:
        return "image"

    if "解説文問題" in joined:
        return "description"

    if "道具名問題" in joined:
        return "term"

    return "term"


def cleanup_noise(lines: list[str]) -> list[str]:
    cleaned: list[str] = []

    for line in lines:
        if re.fullmatch(r"\d+\s*/\s*\d+", line):
            continue
        cleaned.append(line)

    return cleaned


def remove_header_noise_from_question(lines: list[str]) -> list[str]:
    filtered: list[str] = []

    for line in lines:
        if Q_RE.fullmatch(line):
            continue
        if (
            "イラスト問題" in line
            or "写真問題" in line
            or "解説文問題" in line
            or "道具名問題" in line
        ):
            continue
        if line == "問題":
            continue
        if line.endswith("クイズ 基本編"):
            continue
        filtered.append(line)

    return filtered


def remove_header_noise_from_explanation(lines: list[str]) -> list[str]:
    filtered: list[str] = []

    for line in lines:
        if Q_RE.fullmatch(line):
            continue
        if "正解・解説" in line:
            continue
        if line.endswith("クイズ 基本編"):
            continue
        filtered.append(line)

    return filtered


def extract_main_picture(slide, q_no: int, slug: str) -> str | None:
    """
    問題スライド内の一番大きい画像を保存する。
    """
    candidates: list[tuple[int, Any]] = []

    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                area = int(shape.width) * int(shape.height)
                candidates.append((area, shape))
        except Exception:
            continue

    if not candidates:
        return None

    candidates.sort(key=lambda x: x[0], reverse=True)
    picture_shape = candidates[0][1]

    try:
        image = picture_shape.image
        ext = image.ext.lower()
        if ext == "jpeg":
            ext = "jpg"
        if ext not in {"png", "jpg", "gif", "bmp", "tiff", "wmf", "emf"}:
            ext = "png"

        out_dir = Path("public/images/construction") / slug
        out_dir.mkdir(parents=True, exist_ok=True)

        filename = f"q{q_no:03d}.{ext}"
        out_path = out_dir / filename
        out_path.write_bytes(image.blob)

        return f"/images/construction/{slug}/{filename}"
    except Exception:
        return None


def parse_question_slide(lines: list[str]) -> dict[str, Any]:
    src = cleanup_noise(lines)
    src = remove_header_noise_from_question(src)

    q_no = find_q_number(lines)
    kind = detect_kind(lines)

    hint = None
    choices: list[str] = []
    question_lines: list[str] = []

    i = 0
    while i < len(src):
        line = src[i]

        if line.startswith("ヒント"):
            hint = line
            i += 1
            continue

        if line in {"A", "B", "C", "D", "Ａ", "Ｂ", "Ｃ", "Ｄ"}:
            if i + 1 < len(src):
                choices.append(normalize_text(src[i + 1]))
                i += 2
                continue

        m = CHOICE_RE.match(line)
        if m:
            choices.append(normalize_text(m.group(1)))
            i += 1
            continue

        question_lines.append(line)
        i += 1

    question = clean_question_text("\n".join(question_lines).strip())
    hint = clean_hint_text(hint)

    return {
        "id": q_no,
        "kind": kind,
        "question": question,
        "choices": choices,
        "hint": hint,
    }


def parse_explanation_slide(lines: list[str]) -> dict[str, Any]:
    src = cleanup_noise(lines)
    src = remove_header_noise_from_explanation(src)

    q_no = find_q_number(lines)

    correct_letter = None
    explanation_lines: list[str] = []

    in_answer = False
    in_explanation = False

    for line in src:
        if "✓ 正解" in line:
            in_answer = True
            in_explanation = False
            continue

        if "📋 解説" in line or line == "解説":
            in_answer = False
            in_explanation = True
            continue

        if in_answer and correct_letter is None:
            if line and line[0] in LETTER_TO_INDEX:
                correct_letter = line[0]

        if in_explanation:
            explanation_lines.append(line)

    explanation = clean_explanation_text("\n".join(explanation_lines).strip())

    return {
        "id": q_no,
        "correctIndex": LETTER_TO_INDEX.get(correct_letter),
        "explanation": explanation,
    }


def merge_question_and_explanation(
    q: dict[str, Any],
    e: dict[str, Any],
) -> dict[str, Any]:
    if q["id"] != e["id"]:
        raise ValueError(f"Q番号不一致: {q['id']} != {e['id']}")

    if not q["question"]:
        raise ValueError(f"Q{q['id']}: 問題文が空")

    if len(q["choices"]) != 4:
        raise ValueError(f"Q{q['id']}: 選択肢が4件ではない -> {q['choices']}")

    if e["correctIndex"] is None:
        raise ValueError(f"Q{q['id']}: 正解を取得できない")

    item: dict[str, Any] = {
        "id": q["id"],
        "sectionId": "all",
        "kind": q["kind"],
        "question": q["question"],
        "choices": q["choices"],
        "correctIndex": e["correctIndex"],
        "explanation": e["explanation"],
    }

    if q.get("hint"):
        item["hint"] = q["hint"]

    if q.get("imageUrl"):
        item["imageUrl"] = q["imageUrl"]

    return item


def slug_from_filename(filename: str) -> str:
    name = Path(filename).stem

    mapping = {
        "建設道具クイズ": "construction-tools",
        "空調衛生用語クイズ": "hvac-terms",
        "プラント用語クイズ": "plant-terms",
        "建築用語クイズ": "architecture-terms",
        "施工管理用語クイズ": "construction-management-terms",
        "電気用語クイズ": "electric-terms",
        "土木用語クイズ": "civil-terms",
    }

    return mapping.get(name, name)


def title_from_slug(slug: str) -> str:
    mapping = {
        "construction-tools": "建設道具クイズ",
        "hvac-terms": "空調衛生用語クイズ",
        "plant-terms": "プラント用語クイズ",
        "architecture-terms": "建築用語クイズ",
        "construction-management-terms": "施工管理用語クイズ",
        "electric-terms": "電気用語クイズ",
        "civil-terms": "土木用語クイズ",
    }
    return mapping.get(slug, slug)


def const_name_from_slug(slug: str) -> str:
    parts = re.split(r"[^a-zA-Z0-9]+", slug)
    parts = [p for p in parts if p]
    if not parts:
        return "quizData"
    return parts[0] + "".join(p[:1].upper() + p[1:] for p in parts[1:])


def to_ts_const(
    items: list[dict[str, Any]],
    const_name: str,
    quiz_id: str,
    title: str,
) -> str:
    json_str = json.dumps(items, ensure_ascii=False, indent=2)

    return (
        'import type { Quiz } from "../types"\n\n'
        f"export const {const_name}Quiz: Quiz = {{\n"
        f'  id: "{quiz_id}",\n'
        f'  title: "{title}",\n'
        f'  description: "{title}",\n'
        f'  sections: [{{ id: "all", label: "すべて" }}],\n'
        f"  questions: {json_str}\n"
        f"}}\n\n"
        f"export default {const_name}Quiz\n"
    )


def main() -> None:
    if len(sys.argv) < 2:
        print('使い方: python scripts/extract_construction_quiz.py "建設道具クイズ.pptx"')
        sys.exit(1)

    ppt_path = Path(sys.argv[1])
    if not ppt_path.exists():
        print(f"ファイルが見つかりません: {ppt_path}")
        sys.exit(1)

    slug = slug_from_filename(ppt_path.name)

    out_dir = Path("app/data/quizzes")
    out_dir.mkdir(parents=True, exist_ok=True)

    out_json = out_dir / f"{slug}.quiz.json"
    out_ts = out_dir / f"{slug}.quiz.ts"

    const_name = const_name_from_slug(slug)
    title = title_from_slug(slug)

    prs = Presentation(str(ppt_path))

    question_map: dict[int, dict[str, Any]] = {}
    explanation_map: dict[int, dict[str, Any]] = {}

    for slide in prs.slides:
        lines = get_slide_lines(slide)
        if not lines:
            continue

        if is_answer_list_slide(lines):
            continue

        q_no = find_q_number(lines)
        if q_no is None:
            continue

        if is_explanation_slide(lines):
            explanation_map[q_no] = parse_explanation_slide(lines)
        else:
            q = parse_question_slide(lines)

            if q["kind"] == "image" and q_no is not None:
                image_url = extract_main_picture(slide, q_no, slug)
                if image_url:
                    q["imageUrl"] = image_url

            question_map[q_no] = q

    all_q_numbers = sorted(set(question_map.keys()) | set(explanation_map.keys()))

    missing_questions = [n for n in all_q_numbers if n not in question_map]
    missing_explanations = [n for n in all_q_numbers if n not in explanation_map]

    if missing_questions:
        raise ValueError(f"問題スライド不足: {missing_questions}")
    if missing_explanations:
        raise ValueError(f"解説スライド不足: {missing_explanations}")

    merged: list[dict[str, Any]] = []
    for q_no in sorted(question_map.keys()):
        merged.append(
            merge_question_and_explanation(
                question_map[q_no],
                explanation_map[q_no],
            )
        )

    out_json.write_text(
        json.dumps(merged, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    out_ts.write_text(
        to_ts_const(merged, const_name, slug, title),
        encoding="utf-8",
    )

    print("✅ 抽出完了")
    print(f"  source: {ppt_path.name}")
    print(f"  questions: {len(merged)}")
    print(f"  json: {out_json}")
    print(f"  ts:   {out_ts}")


if __name__ == "__main__":
    main()